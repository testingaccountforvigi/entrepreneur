"""Generate the redesigned Aranya Carbon pitch deck (content preserved, visuals rebuilt)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import deck_lib as D

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def new_slide():
    return prs.slides.add_slide(BLANK)


# ============================================================ SLIDE 1 - HERO
def slide_title():
    s = new_slide()
    bgs = D.rect(s, -0.06, -0.06, 13.45, 7.62, fill=D.INK)
    D.set_gradient(bgs, [(0, "0A2A1E"), (55, "103A28"), (100, "16532F")], angle=120)
    # faint large watermark icon
    D.icon(s, "tree-pine", 9.7, 1.0, 4.6, "13FF8F", stroke=0.9)  # will dim via alpha overlay
    veil = D.rect(s, 9.6, 0.9, 4.9, 5.4, fill=D.INK)
    D.set_gradient(veil, [(0, "0A2A1E"), (100, "16532F")], angle=120)
    # bring watermark "into" bg by drawing it faint over veil
    D.icon(s, "leaf", 10.9, 4.5, 3.0, "1C6B3C", stroke=1.0)
    # top + bottom gold hairlines
    D.rect(s, 0, 0, 13.333, 0.09, fill=D.GOLD)
    D.rect(s, 0, 7.41, 13.333, 0.09, fill=D.EMERALD)

    # brand emblem (circle ring + leaf)
    D.oval(s, 6.06, 1.35, 1.22, 1.22, fill="103A28", line=D.EMERALD, line_w=1.6)
    D.icon(s, "tree-pine", 6.4, 1.66, 0.56, D.EMERALD, stroke=1.8)

    # wordmark
    D.text(s, "Aranya Carbon", 0, 2.78, 13.333, 1.1, 56, color=D.WHITE, bold=True,
           font=D.DISPLAY, align=PP_ALIGN.CENTER)

    # tagline pill
    pill_w, pill_h = 5.7, 0.62
    px = (13.333 - pill_w) / 2
    pill = D.rrect(s, px, 4.0, pill_w, pill_h, fill="10402B", line=D.EMERALD, line_w=1.3, radius=0.31)
    D.icon(s, "globe", px + 0.34, 4.16, 0.3, D.GOLD, stroke=1.9)
    D.text(s, "The Stripe of Indian community carbon", px + 0.5, 4.0, pill_w - 0.6, pill_h,
           14, color=D.WHITE, bold=True, font=D.BODY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # subtitle
    D.text(s, "Forests store carbon. Communities earn nothing.", 0, 4.95, 13.333, 0.5,
           19, color=D.ONDARK_SUB, font=D.BODY, align=PP_ALIGN.CENTER)

    # footer
    D.text(s, "Team WeTookAi'sJob   ·   Business Model Canvas", 0, 6.92, 13.333, 0.4,
           12, color="9FBCA9", font=D.BODY, align=PP_ALIGN.CENTER, tracking=1.2)


# ============================================================ helpers for cards
def stat_card(s, l, t, w, h, name, title, sub, accent=D.EMERALD, icon_clr=D.FOREST,
              chip=D.SAGE_LT, featured=False):
    card = D.rrect(s, l, t, w, h, fill=D.WHITE, line=D.BORDER, line_w=1.0, radius=0.18)
    D.add_shadow(card, blur=0.11, dist=0.05, alpha=15000)
    # top accent bar
    bar = D.rrect(s, l, t, w, 0.13, fill=accent, radius=0.06)
    D.icon_chip(s, name, l + 0.34, t + 0.42, 0.96, chip_fill=chip, icon_color=icon_clr, radius=0.22)
    D.two_line(s, l + 0.34, t + 1.62, w - 0.68, h - 1.7, title, sub, 15.5, 11.5,
               t_color=D.INK, s_color=D.MUTED, gap=5)
    return card


# ============================================================ SLIDE 2 - SEGMENTS
def slide_segments():
    s = new_slide(); D.bg(s); D.header(s, "WHO WE SERVE", "Customer Segments")
    D.corner_leaf(s)
    data = [
        ("tree-pine", "Forest communities", "Gram panchayats"),
        ("building-2", "Corporates", "With net-zero goals"),
        ("handshake", "NGOs", "& forest departments"),
        ("coins", "Institutional buyers", "Carbon funds"),
    ]
    accents = [D.EMERALD, D.FOREST, D.EMERALD, D.GOLD]
    icols = [D.FOREST, D.FOREST, D.FOREST, D.GOLD]
    chips = [D.SAGE_LT, D.SAGE_LT, D.SAGE_LT, D.GOLD_LT]
    n = 4; gap = 0.36; left = 0.92; right = 12.41
    w = (right - left - gap * (n - 1)) / n
    top, h = 2.25, 3.95
    for i, (ic, tt, su) in enumerate(data):
        stat_card(s, left + i * (w + gap), top, w, h, ic, tt, su,
                  accent=accents[i], icon_clr=icols[i], chip=chips[i])
    D.footer(s, "02")


# ============================================================ SLIDE 3 - VALUE PROP
def slide_value():
    s = new_slide(); D.bg(s); D.header(s, "WHY IT MATTERS", "Value Proposition")
    # problem banner (dark) - exact original text preserved
    ban = D.rrect(s, 0.92, 2.22, 11.49, 0.92, fill=D.INK, radius=0.16)
    D.set_gradient(ban, [(0, "0C2A1F"), (100, "17402B")], angle=0)
    D.add_shadow(ban, blur=0.1, dist=0.05, alpha=18000)
    D.icon_chip(s, "link-2", 1.22, 2.42, 0.52, chip_fill="1C3B2C", icon_color=D.GOLD, radius=0.13)
    D.text_runs(s, [
        {"text": "Problem:  ", "size": 16.5, "color": D.GOLD, "bold": True, "font": D.BODY},
        {"text": "Verification costs too much, takes years", "size": 16.5,
         "color": D.WHITE, "bold": True, "font": D.DISPLAY},
    ], 1.95, 2.22, 10.2, 0.92, anchor=MSO_ANCHOR.MIDDLE)

    # two split cards - preserve title/sub hierarchy and casing exactly
    cw, ch, ct = 5.58, 2.95, 3.45
    # left
    c1 = D.rrect(s, 0.92, ct, cw, ch, fill=D.WHITE, line=D.BORDER, line_w=1.0, radius=0.2)
    D.add_shadow(c1, blur=0.12, dist=0.055, alpha=15000)
    D.rrect(s, 0.92, ct, cw, 0.14, fill=D.EMERALD, radius=0.07)
    D.icon_chip(s, "tree-pine", 1.32, ct + 0.5, 1.05, chip_fill=D.SAGE_LT, icon_color=D.FOREST, radius=0.24)
    D.two_line(s, 1.32, ct + 1.72, cw - 0.8, 1.1, "For forest owners",
               "Direct income, zero upfront cost", 21, 14.5, t_color=D.INK,
               s_color=D.MUTED, gap=7)
    # right
    rx = 6.84
    c2 = D.rrect(s, rx, ct, cw, ch, fill=D.WHITE, line=D.BORDER, line_w=1.0, radius=0.2)
    D.add_shadow(c2, blur=0.12, dist=0.055, alpha=15000)
    D.rrect(s, rx, ct, cw, 0.14, fill=D.GOLD, radius=0.07)
    D.icon_chip(s, "building", rx + 0.4, ct + 0.5, 1.05, chip_fill=D.GOLD_LT, icon_color=D.GOLD, radius=0.24)
    D.two_line(s, rx + 0.4, ct + 1.72, cw - 0.8, 1.1, "For corporate buyers",
               "Traceable Indian community credits", 21, 14.5, t_color=D.INK,
               s_color=D.MUTED, gap=7)
    D.footer(s, "03")


# ============================================================ SLIDE 4 - CHANNELS (timeline)
def slide_channels():
    s = new_slide(); D.bg(s); D.header(s, "HOW WE REACH THEM", "Channels")
    D.corner_leaf(s)
    steps = [
        ("handshake", "NGO partner network"),
        ("footprints", "Field operations team"),
        ("monitor", "B2B credit marketplace"),
        ("bar-chart-3", "Corporate ESG outreach"),
    ]
    n = 4; gap = 0.42; left = 0.92; right = 12.41
    w = (right - left - gap * (n - 1)) / n
    top, h = 2.65, 3.3
    cy = top + 0.95  # connector line height (through icon centers)
    # connecting line behind cards
    D.connector(s, left + w * 0.5, cy, right - w * 0.5, cy, color=D.LINE_FNT, w=2.2)
    for i, (ic, label) in enumerate(steps):
        cx = left + i * (w + gap)
        card = D.rrect(s, cx, top, w, h, fill=D.WHITE, line=D.BORDER, line_w=1.0, radius=0.18)
        D.add_shadow(card, blur=0.11, dist=0.05, alpha=14000)
        # step number badge
        D.text(s, f"0{i+1}", cx + 0.32, top + 0.26, 1.0, 0.4, 12, color=D.GOLD, bold=True,
               font=D.DISPLAY, tracking=1.0)
        # icon circle (on the line)
        chip_sz = 1.18
        D.icon_chip(s, ic, cx + (w - chip_sz) / 2, top + 0.62, chip_sz, chip_fill=D.SAGE_LT,
                    icon_color=D.FOREST, circle=True, line=D.WHITE, shadow=True)
        D.text(s, label, cx + 0.2, top + 2.1, w - 0.4, 1.0, 14, color=D.INK, bold=True,
               font=D.DISPLAY, align=PP_ALIGN.CENTER, line_spacing=1.05)
        # arrow between cards
        if i < n - 1:
            ax = cx + w + gap / 2
            D.text(s, "›", ax - 0.16, cy - 0.32, 0.4, 0.5, 26, color=D.EMERALD,
                   bold=True, font=D.DISPLAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    D.footer(s, "04")


# ============================================================ SLIDE 5 - RELATIONSHIPS (rows)
def slide_relationships():
    s = new_slide(); D.bg(s); D.header(s, "TRUST & ENGAGEMENT", "Customer Relationships")
    D.corner_leaf(s)
    rows = [
        ("globe", "Local-language field support"),
        ("banknote", "Transparent direct payouts"),
        ("satellite-dish", "Live satellite monitoring"),
        ("user-round", "Dedicated buyer accounts"),
    ]
    left, w = 0.92, 11.49
    top, rh, gap = 2.35, 0.95, 0.26
    for i, (ic, label) in enumerate(rows):
        y = top + i * (rh + gap)
        card = D.rrect(s, left, y, w, rh, fill=D.WHITE, line=D.BORDER, line_w=1.0, radius=0.16)
        D.add_shadow(card, blur=0.09, dist=0.04, alpha=12000)
        # left accent edge
        D.rrect(s, left, y, 0.13, rh, fill=D.EMERALD, radius=0.06)
        D.icon_chip(s, ic, left + 0.34, y + (rh - 0.62) / 2, 0.62, chip_fill=D.SAGE_LT,
                    icon_color=D.FOREST, radius=0.16)
        D.text(s, label, left + 1.2, y, w - 2.2, rh, 16, color=D.TEXT, bold=True,
               font=D.DISPLAY, anchor=MSO_ANCHOR.MIDDLE)
        # trust check on right
        D.icon(s, "badge-check", left + w - 0.78, y + (rh - 0.42) / 2, 0.42, D.SAGE)
    D.footer(s, "05")


# ============================================================ SLIDE 6 - REVENUE (SaaS cards)
def slide_revenue():
    s = new_slide(); D.bg(s); D.header(s, "HOW WE EARN", "Revenue Streams")
    D.corner_leaf(s)
    data = [
        ("circle-dollar-sign", "Credit sale commission", "18% primary", True),
        ("satellite-dish", "Monitoring subscriptions", "Annual per project", False),
        ("puzzle", "NGO SaaS subscriptions", "Monthly platform fee", False),
        ("building", "Enterprise buyer subs", "Priority access", False),
    ]
    n = 4; gap = 0.36; left = 0.92; right = 12.41
    w = (right - left - gap * (n - 1)) / n
    top, h = 2.25, 4.0
    for i, (ic, tt, tag, feat) in enumerate(data):
        cx = left + i * (w + gap)
        if feat:
            card = D.rrect(s, cx, top - 0.12, w, h + 0.24, fill=D.INK, radius=0.2)
            D.set_gradient(card, [(0, "103A28"), (100, "17532F")], angle=120)
            D.add_shadow(card, blur=0.14, dist=0.07, alpha=24000)
            chip_f, icon_c, t_c, tag_bg, tag_c = "1C4632", D.GOLD, D.WHITE, D.GOLD, D.INK
            D.rrect(s, cx, top - 0.12, w, 0.14, fill=D.GOLD, radius=0.07)
            iy = top + 0.42
        else:
            card = D.rrect(s, cx, top, w, h, fill=D.WHITE, line=D.BORDER, line_w=1.0, radius=0.2)
            D.add_shadow(card, blur=0.11, dist=0.05, alpha=14000)
            chip_f, icon_c, t_c, tag_bg, tag_c = D.SAGE_LT, D.FOREST, D.INK, D.SAGE_LT2, D.EMERALD_D
            iy = top + 0.42
        D.icon_chip(s, ic, cx + 0.34, iy, 0.96, chip_fill=chip_f, icon_color=icon_c, radius=0.22)
        D.text(s, tt, cx + 0.34, iy + 1.18, w - 0.68, 1.1, 15, color=t_c, bold=True,
               font=D.DISPLAY, line_spacing=1.05)
        # tag pill at bottom
        tw = w - 0.68
        D.rrect(s, cx + 0.34, top + h - (0.7 if not feat else 0.58), tw, 0.46,
                fill=tag_bg, radius=0.23)
        D.text(s, tag, cx + 0.34, top + h - (0.7 if not feat else 0.58), tw, 0.46,
               11.5, color=tag_c, bold=True, font=D.BODY, align=PP_ALIGN.CENTER,
               anchor=MSO_ANCHOR.MIDDLE)
    D.footer(s, "06")


# ============================================================ SLIDE 7 - KEY RESOURCES (2x2 stack)
def slide_resources():
    s = new_slide(); D.bg(s); D.header(s, "WHAT WE NEED", "Key Resources")
    D.corner_leaf(s)
    data = [
        ("monitor", "Digital platform & app"),
        ("satellite-dish", "Satellite & geospatial data"),
        ("cpu", "Carbon data & AI models"),
        ("users", "Field team & domain experts"),
    ]
    left, top = 0.92, 2.3
    gapx, gapy = 0.36, 0.32
    w = (12.41 - left - gapx) / 2
    h = 1.85
    for i, (ic, tt) in enumerate(data):
        cx = left + (i % 2) * (w + gapx)
        cy = top + (i // 2) * (h + gapy)
        card = D.rrect(s, cx, cy, w, h, fill=D.WHITE, line=D.BORDER, line_w=1.0, radius=0.18)
        D.add_shadow(card, blur=0.11, dist=0.05, alpha=14000)
        D.rrect(s, cx, cy, 0.13, h, fill=D.EMERALD, radius=0.06)
        D.icon_chip(s, ic, cx + 0.42, cy + (h - 1.05) / 2, 1.05, chip_fill=D.SAGE_LT,
                    icon_color=D.FOREST, radius=0.22)
        D.text(s, tt, cx + 1.72, cy, w - 2.0, h, 18, color=D.INK, bold=True,
               font=D.DISPLAY, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.04)
    D.footer(s, "07")


# ============================================================ SLIDE 8 - KEY ACTIVITIES (cycle)
def slide_activities():
    s = new_slide(); D.bg(s); D.header(s, "WHAT WE DO", "Key Activities")
    cx, cy = 6.666, 4.35   # center of cycle
    ring_r = 2.05
    # ring
    D.oval(s, cx - ring_r, cy - ring_r, ring_r * 2, ring_r * 2, fill=None,
           line=D.LINE_FNT, line_w=2.0)
    # center hub
    hub = 1.5
    hubsh = D.oval(s, cx - hub / 2, cy - hub / 2, hub, hub, fill=D.INK)
    D.set_gradient(hubsh, [(0, "12462C"), (100, "0C2A1F")], angle=90)
    D.add_shadow(hubsh, blur=0.12, dist=0.05, alpha=22000)
    D.icon(s, "repeat", cx - 0.34, cy - 0.5, 0.68, D.GOLD, stroke=1.9)
    D.text(s, "Continuous", cx - 0.9, cy + 0.18, 1.8, 0.3, 11, color=D.WHITE, bold=True,
           font=D.DISPLAY, align=PP_ALIGN.CENTER)
    D.text(s, "cycle", cx - 0.9, cy + 0.42, 1.8, 0.3, 11, color=D.ONDARK_SUB,
           font=D.BODY, align=PP_ALIGN.CENTER)

    nodes = [
        ("map-pin", "Community onboarding & mapping", "1"),
        ("tree-pine", "Carbon measurement & PDD", "2"),
        ("satellite-dish", "Monitoring & verification", "3"),
        ("banknote", "Marketplace sales & payouts", "4"),
    ]
    # positions: top, right, bottom, left
    pos = [(cx, cy - ring_r), (cx + ring_r, cy), (cx, cy + ring_r), (cx - ring_r, cy)]
    cardw, cardh = 3.15, 1.18
    placements = [
        (cx - cardw / 2, cy - ring_r - cardh - 0.12),      # top
        (cx + ring_r + 0.18, cy - cardh / 2),              # right
        (cx - cardw / 2, cy + ring_r + 0.12),              # bottom
        (cx - ring_r - cardw - 0.18, cy - cardh / 2),      # left
    ]
    for i, (ic, label, num) in enumerate(nodes):
        l, t = placements[i]
        card = D.rrect(s, l, t, cardw, cardh, fill=D.WHITE, line=D.BORDER, line_w=1.0, radius=0.16)
        D.add_shadow(card, blur=0.1, dist=0.045, alpha=15000)
        D.icon_chip(s, ic, l + 0.22, t + (cardh - 0.74) / 2, 0.74, chip_fill=D.SAGE_LT,
                    icon_color=D.FOREST, radius=0.18)
        D.text(s, label, l + 1.08, t, cardw - 1.24, cardh, 12.5, color=D.INK, bold=True,
               font=D.DISPLAY, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        # step number badge on ring point
        bx, by = pos[i]
        D.oval(s, bx - 0.2, by - 0.2, 0.4, 0.4, fill=D.GOLD)
        D.text(s, num, bx - 0.2, by - 0.21, 0.4, 0.42, 13, color=D.INK, bold=True,
               font=D.DISPLAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # bottom caption
    cap = D.rrect(s, 2.9, 6.62, 7.53, 0.5, fill=D.SAGE_LT2, radius=0.25)
    D.text(s, "A continuous cycle  ·  onboard → measure → monitor → sell → repeat",
           2.9, 6.62, 7.53, 0.5, 12, color=D.FOREST, bold=True, font=D.BODY,
           align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    D.footer(s, "08")


# ============================================================ SLIDE 9 - KEY PARTNERS (network)
def slide_partners():
    s = new_slide(); D.bg(s); D.header(s, "WHO WE WORK WITH", "Key Partners")
    cx, cy = 6.666, 4.45
    nodes = [
        ("handshake", "NGOs & communities", "Field relationships"),
        ("badge-check", "Third-party verifiers", "Audit & validation"),
        ("clipboard-list", "Registries", "Verra · Gold Standard · BEE"),
        ("credit-card", "Payment & satellite", "Providers & data"),
    ]
    cardw, cardh = 3.25, 1.5
    placements = [
        (0.92, 2.5), (12.41 - cardw, 2.5),
        (0.92, 5.05), (12.41 - cardw, 5.05),
    ]
    centers = [(p[0] + cardw / 2, p[1] + cardh / 2) for p in placements]
    # connectors from hub to each node center (drawn first, behind)
    for (nx, ny) in centers:
        D.connector(s, cx, cy, nx, ny, color=D.LINE_FNT, w=1.8, dash="dash")
    # hub
    hub = 1.85
    hubsh = D.oval(s, cx - hub / 2, cy - hub / 2, hub, hub, fill=D.INK)
    D.set_gradient(hubsh, [(0, "12462C"), (100, "0C2A1F")], angle=90)
    D.add_shadow(hubsh, blur=0.13, dist=0.06, alpha=24000)
    D.oval(s, cx - hub / 2 - 0.08, cy - hub / 2 - 0.08, hub + 0.16, hub + 0.16, fill=None,
           line=D.EMERALD, line_w=1.4)
    D.icon(s, "leaf", cx - 0.42, cy - 0.52, 0.5, D.EMERALD, stroke=1.9)
    D.text(s, "Aranya", cx - 0.9, cy + 0.05, 1.8, 0.3, 13, color=D.WHITE, bold=True,
           font=D.DISPLAY, align=PP_ALIGN.CENTER)
    D.text(s, "Carbon", cx - 0.9, cy + 0.32, 1.8, 0.3, 13, color=D.WHITE, bold=True,
           font=D.DISPLAY, align=PP_ALIGN.CENTER)
    # nodes
    for i, (ic, tt, su) in enumerate(nodes):
        l, t = placements[i]
        card = D.rrect(s, l, t, cardw, cardh, fill=D.WHITE, line=D.BORDER, line_w=1.0, radius=0.18)
        D.add_shadow(card, blur=0.11, dist=0.05, alpha=15000)
        D.icon_chip(s, ic, l + 0.3, t + (cardh - 0.92) / 2, 0.92, chip_fill=D.SAGE_LT,
                    icon_color=D.FOREST, radius=0.2)
        D.two_line(s, l + 1.42, t + 0.34, cardw - 1.6, cardh - 0.5, tt, su, 15, 11,
                   t_color=D.INK, s_color=D.MUTED, gap=4)
    D.footer(s, "09")


# ============================================================ SLIDE 10 - COST STRUCTURE
def slide_cost():
    s = new_slide(); D.bg(s); D.header(s, "WHAT IT COSTS", "Cost Structure")
    D.corner_leaf(s)
    rows = [
        ("footprints", "Field operations"),
        ("monitor", "Technology & infrastructure"),
        ("users", "Salaries"),
        ("megaphone", "Sales & marketing"),
        ("scale", "Admin & legal"),
    ]
    left, w = 0.92, 11.49
    top, rh, gap = 2.3, 0.78, 0.2
    for i, (ic, label) in enumerate(rows):
        y = top + i * (rh + gap)
        card = D.rrect(s, left, y, w, rh, fill=D.WHITE, line=D.BORDER, line_w=1.0, radius=0.15)
        D.add_shadow(card, blur=0.09, dist=0.04, alpha=12000)
        D.rrect(s, left, y, 0.13, rh, fill=D.GOLD, radius=0.06)
        D.icon_chip(s, ic, left + 0.32, y + (rh - 0.52) / 2, 0.52, chip_fill=D.SAGE_LT,
                    icon_color=D.FOREST, radius=0.14)
        D.text(s, label, left + 1.08, y, w - 2.4, rh, 15.5, color=D.TEXT, bold=True,
               font=D.DISPLAY, anchor=MSO_ANCHOR.MIDDLE)
        # decorative cost marker (visual only, no data/words)
        for d in range(3):
            D.oval(s, left + w - 0.95 + d * 0.26, y + rh / 2 - 0.05, 0.1, 0.1,
                   fill=D.GOLD if d == 0 else D.SAGE_LT)
    D.footer(s, "10")


# ============================================================ build
slide_title()
slide_segments()
slide_value()
slide_channels()
slide_relationships()
slide_revenue()
slide_resources()
slide_activities()
slide_partners()
slide_cost()

prs.save("Aranya_Carbon_Pitch_Deck.pptx")
print(f"Saved. Slides: {len(prs.slides)}")
