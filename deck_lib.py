"""Shared design system: palette, typography, vector-icon rendering, and shape helpers."""
import os, re, hashlib
import cairosvg
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- palette
BG       = "F4F7F2"   # off-white neutral
BG2      = "FFFFFF"
INK      = "0C2A1F"   # deep forest, near-black
FOREST   = "14532D"   # primary deep forest green
FOREST2  = "1C6B3C"   # gradient companion
EMERALD  = "12A06A"   # accent emerald
EMERALD_D= "0C7F54"
SAGE_LT  = "E7F0E8"   # icon chip background
SAGE_LT2 = "EFF5EE"
SAGE     = "8FA99A"   # secondary sage green
GOLD     = "C2912E"   # warm gold highlight
GOLD_LT  = "F4E9CB"
WHITE    = "FFFFFF"
TEXT     = "1B2B23"
MUTED    = "5E7167"
BORDER   = "E0E8DE"
LINE_FNT = "D7E1D5"
ONDARK   = "FFFFFF"
ONDARK_SUB = "C9DCCE"

DISPLAY = "Poppins"      # modern geometric headings
BODY    = "Inter"        # clean body
DISPLAY_FB = "Segoe UI"  # graceful fallback handled by PowerPoint automatically

ICON_DIR = "icons_svg"
PNG_DIR  = "icons_png"
os.makedirs(PNG_DIR, exist_ok=True)


def render_icon(name, color, stroke=1.9, px=520):
    """Recolor a Lucide SVG and rasterize to a crisp PNG. Returns path."""
    key = hashlib.md5(f"{name}{color}{stroke}{px}".encode()).hexdigest()[:10]
    out = os.path.join(PNG_DIR, f"{name}_{key}.png")
    if os.path.exists(out):
        return out
    raw = open(os.path.join(ICON_DIR, f"{name}.svg")).read()
    raw = raw.replace("currentColor", f"#{color}")
    raw = re.sub(r'stroke-width="[^"]*"', f'stroke-width="{stroke}"', raw)
    cairosvg.svg2png(bytestring=raw.encode(), write_to=out,
                     output_width=px, output_height=px)
    return out


# ---------------------------------------------------------------- xml helpers
def _spPr(shape):
    return shape._element.spPr


def add_shadow(shape, blur=0.10, dist=0.045, direction=5400000, color=INK, alpha=20000):
    """Subtle soft drop shadow. blur/dist in inches, alpha = shadow opacity (0-100000)."""
    sp = _spPr(shape)
    for el in sp.findall(qn('a:effectLst')):
        sp.remove(el)
    eff = sp.makeelement(qn('a:effectLst'), {})
    sh = sp.makeelement(qn('a:outerShdw'), {
        'blurRad': str(int(Inches(blur))),
        'dist': str(int(Inches(dist))),
        'dir': str(direction), 'rotWithShape': '0'})
    clr = sp.makeelement(qn('a:srgbClr'), {'val': color})
    al = sp.makeelement(qn('a:alpha'), {'val': str(alpha)})
    clr.append(al); sh.append(clr); eff.append(sh)
    sp.append(eff)


def set_gradient(shape, stops, angle=90):
    """stops: list of (pos_0_100, hexcolor). Linear gradient."""
    sp = _spPr(shape)
    for tag in ('a:noFill', 'a:solidFill', 'a:gradFill', 'a:blipFill', 'a:pattFill'):
        for el in sp.findall(qn(tag)):
            sp.remove(el)
    grad = sp.makeelement(qn('a:gradFill'), {})
    lst = sp.makeelement(qn('a:gsLst'), {})
    for pos, hexc in stops:
        gs = sp.makeelement(qn('a:gs'), {'pos': str(int(pos * 1000))})
        c = sp.makeelement(qn('a:srgbClr'), {'val': hexc})
        gs.append(c); lst.append(gs)
    grad.append(lst)
    lin = sp.makeelement(qn('a:lin'), {'ang': str(int(angle * 60000)), 'scaled': '1'})
    grad.append(lin)
    ln = sp.find(qn('a:ln'))
    eff = sp.find(qn('a:effectLst'))
    if ln is not None:
        ln.addprevious(grad)
    elif eff is not None:
        eff.addprevious(grad)
    else:
        sp.append(grad)


def no_line(shape):
    shape.line.fill.background()


# ---------------------------------------------------------------- shapes
def rect(slide, l, t, w, h, fill=None, line=None, line_w=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = RGBColor.from_string(line); s.line.width = Pt(line_w)
    return s


def rrect(slide, l, t, w, h, fill=None, line=None, line_w=1.0, radius=0.14):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.shadow.inherit = False
    try:
        s.adjustments[0] = min(0.5, radius / min(w, h))
    except Exception:
        pass
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = RGBColor.from_string(line); s.line.width = Pt(line_w)
    return s


def oval(slide, l, t, w, h, fill=None, line=None, line_w=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(w), Inches(h))
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = RGBColor.from_string(line); s.line.width = Pt(line_w)
    return s


def connector(slide, x1, y1, x2, y2, color=LINE_FNT, w=1.4, dash=None):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = RGBColor.from_string(color)
    c.line.width = Pt(w)
    if dash:
        ln = c.line._get_or_add_ln()
        d = ln.makeelement(qn('a:prstDash'), {'val': dash})
        ln.append(d)
    c.shadow.inherit = False
    return c


def set_tracking(run, pts):
    rPr = run._r.get_or_add_rPr()
    rPr.set('spc', str(int(pts * 100)))


def text(slide, s, l, t, w, h, size, color=TEXT, bold=False, font=BODY,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, tracking=None, line_spacing=None,
         italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    r = p.add_run(); r.text = s
    f = r.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.name = font; f.color.rgb = RGBColor.from_string(color)
    if tracking:
        set_tracking(r, tracking)
    return tb


def two_line(slide, l, t, w, h, title, sub, t_size, s_size, t_color=INK,
             s_color=MUTED, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             t_font=DISPLAY, s_font=BODY, gap=6):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = title
    r.font.size = Pt(t_size); r.font.bold = True; r.font.name = t_font
    r.font.color.rgb = RGBColor.from_string(t_color)
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = align
        p2.space_before = Pt(gap)
        r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(s_size); r2.font.bold = False; r2.font.name = s_font
        r2.font.color.rgb = RGBColor.from_string(s_color)
    return tb


def text_runs(slide, runs, l, t, w, h, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """runs: list of dicts {text,size,color,bold,font,tracking}."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    for rr in runs:
        r = p.add_run(); r.text = rr["text"]
        f = r.font
        f.size = Pt(rr.get("size", 14)); f.bold = rr.get("bold", False)
        f.name = rr.get("font", BODY)
        f.color.rgb = RGBColor.from_string(rr.get("color", TEXT))
        if rr.get("tracking"):
            set_tracking(r, rr["tracking"])
    return tb


def icon(slide, name, l, t, sz, color, stroke=1.9):
    p = render_icon(name, color, stroke=stroke)
    return slide.shapes.add_picture(p, Inches(l), Inches(t), Inches(sz), Inches(sz))


def icon_chip(slide, name, l, t, sz, chip_fill=SAGE_LT, icon_color=FOREST,
              radius=0.2, icon_scale=0.54, stroke=1.9, line=None, circle=False, shadow=False):
    if circle:
        chip = oval(slide, l, t, sz, sz, fill=chip_fill, line=line, line_w=1.2)
    else:
        chip = rrect(slide, l, t, sz, sz, fill=chip_fill, line=line, line_w=1.2, radius=radius)
    if shadow:
        add_shadow(chip, blur=0.08, dist=0.03, alpha=16000)
    iw = sz * icon_scale
    icon(slide, name, l + (sz - iw) / 2, t + (sz - iw) / 2, iw, icon_color, stroke=stroke)
    return chip


# ---------------------------------------------------------------- slide chrome
def bg(slide, color=BG):
    rect(slide, -0.06, -0.06, 13.45, 7.62, fill=color)


def header(slide, eyebrow, title, divider=GOLD):
    # top hairline accent
    rect(slide, 0, 0, 13.333, 0.085, fill=FOREST)
    rect(slide, 0, 0, 4.4, 0.085, fill=EMERALD)
    text(slide, eyebrow, 0.92, 0.62, 11, 0.34, 11, color=EMERALD_D, bold=True,
         font=BODY, tracking=2.4)
    text(slide, title, 0.9, 0.95, 11.6, 0.82, 33, color=INK, bold=True, font=DISPLAY)
    rrect(slide, 0.93, 1.78, 0.66, 0.07, fill=GOLD, radius=0.035)


def footer(slide, page):
    text(slide, "Aranya Carbon", 0.92, 7.04, 4.0, 0.32, 9.5, color=FOREST, bold=True, font=BODY)
    text(slide, "Team WeTookAi'sJob", 2.55, 7.04, 4.0, 0.32, 9.5, color=MUTED, font=BODY)
    # divider dot
    text(slide, "·", 2.32, 7.02, 0.2, 0.32, 9.5, color=SAGE, font=BODY)
    text(slide, page, 12.0, 7.0, 0.95, 0.34, 10, color=EMERALD_D, bold=True,
         font=BODY, align=PP_ALIGN.RIGHT)


def corner_leaf(slide, color=SAGE_LT, l=11.7, t=5.95, sz=2.4):
    """Faint brand watermark in bottom-right."""
    icon(slide, "leaf", l, t, sz, color, stroke=1.2)
